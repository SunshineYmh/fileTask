import json
import shutil
from django.http import HttpResponse
import os
# from PyPDF2 import PdfReader
from django.core.cache import cache
import fitz
from pptx import Presentation
from pptx.util import Inches
from PDFesc.models.fileModels import FileStorage
from .utils import isFIleExists, unique_id, UPLOAD_FILE_PPTX_PATH, DOWNLOAD_FILE_PPTX_PATH

# https://dothinking.github.io/pdf2docx/quickstart.convert.html


def img2pptx(img_path, ppt_name):
    ppt = Presentation()  # 生成ppt对象
    for i in img_path:
        layout = ppt.slide_layouts[6]  # 定义一个 PPT 页面 插入图片,这里选择样式6
        slide = ppt.slides.add_slide(layout)

        image = slide.shapes.add_picture(image_file=i,
                                         left=Inches(0),
                                         top=Inches(0),
                                         width=Inches(10),
                                         height=Inches(8))
        # 目前设置会撑满整张PPT，大小可自行调整
    ppt.save(ppt_name)


def pdf2img(pdf_path, img_dir, ppt_name):
    doc = fitz.open(pdf_path)  # 打开pdf
    img_path = []
    for page in doc:  # 遍历pdf的每一页
        zoom_x = 2.0  # 设置每页的水平缩放因子
        zoom_y = 2.0  # 设置每页的垂直缩放因子
        mat = fitz.Matrix(zoom_x, zoom_y)
        pix = page.get_pixmap(matrix=mat)
        pix.save(r"{}page-{}.png".format(img_dir + "\\", page.number))  # 保存
        img_path.append(r"{}page-{}.png".format(img_dir + "\\", page.number))
    img2pptx(img_path, ppt_name)


def PyPDFToPptx(pdf_file, fileName):
    print(DOWNLOAD_FILE_PPTX_PATH, "--333---", pdf_file, fileName)
    isFIleExists(DOWNLOAD_FILE_PPTX_PATH)
    fileName = fileName.replace('.pdf', '').replace('.PDF', '')
    img_dir = os.path.join(DOWNLOAD_FILE_PPTX_PATH, fileName)
    isFIleExists(img_dir)
    pptx_file = os.path.join(DOWNLOAD_FILE_PPTX_PATH, fileName + '.pptx')

    pdf2img(pdf_file, img_dir, pptx_file)

    file_size = os.stat(pptx_file).st_size
    unique_id_d = unique_id()
    fileName = fileName + '.pptx'
    FileStorage.objects.create(
        file_id=unique_id_d,
        file_name=fileName,
        file_path=pptx_file,
        file_content_type=
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        file_ywfl="pdfToPptx",
        file_ywdl="downloadfile",
        file_size=file_size)
    shutil.rmtree(img_dir)
    context = {unique_id_d: fileName}
    return context


# 上传文件
def upload_file(request):
    if request.method == "POST":
        # newfile = request.FILES.get('newfile', None)
        newfile = request.FILES.getlist('newfile')  # 获得多个文件上传进来的文件列表。
        if not newfile:
            context = {"succse": False, "msg": '提交无效，没有文件上传！'}
            return HttpResponse(json.dumps(context))
        isFIleExists(UPLOAD_FILE_PPTX_PATH)
        down_file_list = {}
        for fileds in newfile:
            upload_file_up = os.path.join(UPLOAD_FILE_PPTX_PATH, fileds.name)
            down_file_list[upload_file_up] = fileds.name
            to_path = open(upload_file_up, 'wb+')
            for chunk in fileds.chunks():
                to_path.write(chunk)

            unique_id_d = unique_id()
            # 设置缓存，设置失效时间为1小时
            # cache.set(unique_id_d, fileds.name, timeout=10)
            FileStorage.objects.create(file_id=unique_id_d,
                                       file_name=fileds.name,
                                       file_path=upload_file_up,
                                       file_content_type="application/pdf",
                                       file_ywfl="pdfToWord",
                                       file_ywdl="uploadfile",
                                       file_size=fileds.size)
        to_path.close()

        docx_data_list = []
        # 遍历字典的键值对
        for downfilePath, filee_name in down_file_list.items():
            docx_data = PyPDFToPptx(downfilePath, filee_name)
            docx_data_list.append(docx_data)

        context = {"succse": True, "data": docx_data_list}
        return HttpResponse(json.dumps(context))
    else:
        context = {"succse": False, "msg": '非表单提交访问！'}
        return HttpResponse(json.dumps(context))
